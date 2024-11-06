from pptx import Presentation
from pptx.util import Inches

# Create a new presentation
prs = Presentation()

# Slide 1: Title Slide
slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Artificial Intelligence and Machine Learning"
subtitle = slide.placeholders[1]
subtitle.text = "Introduction to Group Members and Topic Overview"

# Slide 2: Introduction to AI and ML
slide_layout = prs.slide_layouts[1]
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Introduction to AI and ML"
content = slide.placeholders[1]
content.text = "Overview of AI and ML, history, and evolution."

# Slide 3: Types of AI
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Types of AI"
content = slide.placeholders[1]
content.text = "Narrow AI, General AI, Super intelligent AI. Emphasize Narrow AI's importance."

# Slide 4: Types of Machine Learning
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Types of Machine Learning"
content = slide.placeholders[1]
content.text = "Supervised, Unsupervised, and Reinforcement Learning with examples."

# Slide 5: Key Concepts in ML
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Key Concepts in ML"
content = slide.placeholders[1]
content.text = "Datasets, features, labels, model training/testing, algorithms."

# Slide 6: Applications of AI and ML
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Applications of AI and ML"
content = slide.placeholders[1]
content.text = "Industries: healthcare, finance, automotive, e-commerce."

# Slide 7: Ethical Concerns in AI
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Ethical Concerns in AI"
content = slide.placeholders[1]
content.text = "Bias, privacy, job displacement. Solutions: transparency and fairness."

# Slide 8: Future of AI and ML
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Future of AI and ML"
content = slide.placeholders[1]
content.text = "Emerging trends, anticipated advancements, responsible innovation."

# Slide 9: Case Study
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Case Study: AlphaGo"
content = slide.placeholders[1]
content.text = "How AlphaGo works, its impact, public reaction."

# Slide 10: Conclusion
slide = prs.slides.add_slide(slide_layout)
title = slide.shapes.title
title.text = "Conclusion"
content = slide.placeholders[1]
content.text = "Summary of key points."

# Save the presentation
prs.save('AI_ML_Presentation.pptx')
